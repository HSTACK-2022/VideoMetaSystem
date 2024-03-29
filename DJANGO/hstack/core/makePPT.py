import os
import platform
import tempfile

from . import models
from pptx import Presentation
from pptx.util import Inches


# 상수 설정
OS = platform.system()

def getPPTFile(videoId):
    pptFile = Presentation()
    videoPathObj = models.Videopath.objects.get(id=videoId)     #DB에서 videoId에 해당하는 객체를 가져옴
    imageFilePath = videoPathObj.imageaddr                      #DB에서 imageAddr 추출
    title = videoPathObj.title

    if OS == "Windows" : 
        imagePath = (imageFilePath + "\\").replace("/", "\\")
    else : 
        imagePath = (imageFilePath + "\\").replace("\\", "/")
    
    imageList = os.listdir(imageFilePath)

    slide_layout = pptFile.slide_layouts[6]

    for image in imageList:
        if image.startswith("L") or image.startswith("P"):
            print("**")
            slide = pptFile.slides.add_slide(slide_layout)
            slide.shapes.add_picture(imagePath + image, 0, 0, pptFile.slide_width, pptFile.slide_height)
    
    if OS == "Windows": 
        pptPathRel = "../../../media/Uploaded/Image/" + imageFilePath.split('Image\\')[1].replace("\\", "/") + "/" + title + '.pptx'
    else: 
        pptPathRel = "../../../media/Uploaded/Image/" + imageFilePath.split('Image/')[1] + "/" + title + ".pptx"

    pptFile.save(imagePath + title + '.pptx')

    return pptPathRel

# idea from PR82 by @minzix
# PDF로 저장
from img2pdf import convert

def makePDFFile(videoId): 
    videoPathObj = models.Videopath.objects.get(id=videoId)     #DB에서 videoId에 해당하는 객체를 가져옴
    imageFilePath = videoPathObj.imageaddr                      #DB에서 imageAddr 추출
    title = videoPathObj.title

    if OS == "Windows" : 
        imagePath = (imageFilePath + "\\").replace("/", "\\")
    else : 
        imagePath = (imageFilePath + "\\").replace("\\", "/")

    imageList = os.listdir(imageFilePath)
   
    if OS == "Windows": 
        pdfPathRel = "../../../media/Uploaded/Image/" + imageFilePath.split('Image\\')[1].replace("\\", "/") + "/" + title + '.pdf'
    else: 
        pdfPathRel = "../../../media/Uploaded/Image/" + imageFilePath.split('Image/')[1] + "/" + title + ".pdf"

    imageList = os.listdir(imagePath)
    
    with open(pdfPathRel, "wb") as pf:
        pdfList = []

        for image in imageList:
            if image.startswith("L") or image.startswith("P"):
                pdfList.append(imagePath + image)
        
        pdfFile = convert(pdfList)
        pf.write(pdfFile)