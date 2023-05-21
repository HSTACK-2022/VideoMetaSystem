# makePPT.py
#
# image 탭에 띄울 사진과 PPT 파일을 생성 및 반환하는 코드입니다.
# 
# uses
# - getPPTImage(fileURL) : 이미지 형식이 L(Lecture), P(PPT)인 파일들을 불러옵니다.
# - getPPTFile(fileURL, title) : PPT 파일을 저장합니다.
#
# parameters
# - fileURL : 영상이 저장된 경로
# - title : PPT 파일의 이름


import os

from hstack.config import OS
from pptx import Presentation


# 화면에 띄우기 위한 PPT 이미지 얻기
def getPPTImage(fileURL):
    pptImage = set()
    imagePath = os.path.join(os.path.dirname(fileURL), 'Image')
    imageList = os.listdir(imagePath)

    for image in imageList:
        if image.startswith("L") or image.startswith("P"):
            imageSrc = os.path.join(imagePath, image)
            pptImage.add(imageSrc)
    
    return pptImage


# PPT 파일 얻기
def getPPTFile(fileURL, title):
    pptFile = Presentation()
    imagePath = os.path.join(os.path.dirname(fileURL), 'Image')

    imageList = os.listdir(imagePath)

    slide_layout = pptFile.slide_layouts[6]

    for image in imageList:
        if image.startswith("L") or image.startswith("P"):
            path = os.path.join(imagePath, image)
            slide = pptFile.slides.add_slide(slide_layout)
            slide.shapes.add_picture(path, 0, 0, pptFile.slide_width, pptFile.slide_height)
        
    pptPathRel = os.path.join(os.path.dirname(fileURL), title + '.ppt')
    pptFile.save(pptPathRel)


# idea from PR82 by @minzix

# PDF로 저장
from img2pdf import convert

def makePDFFile(fileURL, title): 
    pdfPathRel = os.path.join(os.path.dirname(fileURL), title + '.pdf')
    imagePath = os.path.join(os.path.dirname(fileURL), 'Image')

    imageList = os.listdir(imagePath)
    
    with open(pdfPathRel, "wb") as pf:
        pdfList = []

        for image in imageList:
            if image.startswith("L") or image.startswith("P"):
                path = os.path.join(imagePath, image)
                pdfList.append(path)
        
        pdfFile = convert(pdfList)
        pf.write(pdfFile)
